from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Cm, Pt


PROJECT = Path(__file__).resolve().parents[1]
TEMPLATE = PROJECT / "sources" / "交大模版.pptx"
OUT = PROJECT / "exports" / "trial_lecture_ai_research_methods_sjtu_visual_full_20260614.pptx"

SJTU_RED = RGBColor(198, 18, 35)
DARK = RGBColor(43, 52, 68)
MUTED = RGBColor(92, 99, 112)
LIGHT_BG = RGBColor(248, 249, 251)
LINE = RGBColor(218, 222, 228)
WHITE = RGBColor(255, 255, 255)
FONT = "Microsoft YaHei"
CONTENT_X = 1.35
CONTENT_W = 30.6
COL_GAP = 0.9
COL_W = (CONTENT_W - COL_GAP * 2) / 3
COL_XS = [CONTENT_X, CONTENT_X + COL_W + COL_GAP, CONTENT_X + (COL_W + COL_GAP) * 2]


def clear_slides(prs):
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        sld_id_lst.remove(sld_id)


def set_font(run, size, color=DARK, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    for name in ("latin", "ea", "cs"):
        el = r_pr.find(qn(f"a:{name}"))
        if el is None:
            el = OxmlElement(f"a:{name}")
            r_pr.append(el)
        el.set("typeface", FONT)


def add_textbox(slide, x, y, w, h, text, size=18, color=DARK, bold=False,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.08,
                line_spacing=1.18):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Cm(margin)
    tf.margin_right = Cm(margin)
    tf.margin_top = Cm(margin)
    tf.margin_bottom = Cm(margin)
    tf.vertical_anchor = valign
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ""
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        set_font(run, size, color=color, bold=bold)
    return box


def set_run_font_family(run):
    run.font.name = FONT
    r_pr = run._r.get_or_add_rPr()
    for name in ("latin", "ea", "cs"):
        el = r_pr.find(qn(f"a:{name}"))
        if el is None:
            el = OxmlElement(f"a:{name}")
            r_pr.append(el)
        el.set("typeface", FONT)


def fill_template_placeholder(slide, placeholder_idx, text):
    ph = next(
        shape for shape in slide.placeholders
        if shape.placeholder_format.idx == placeholder_idx
    )
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    # Keep size, color, position, and alignment inherited from the template.
    set_run_font_family(run)
    return ph


def add_title_bar(slide, title, page_no):
    # Body pages use the template's built-in red masthead and title placeholder.
    fill_template_placeholder(slide, 14, title)


def add_footer(slide, page_no=None):
    return None


def add_card(slide, x, y, w, h, title, body, accent=False):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255) if not accent else LIGHT_BG
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    add_textbox(slide, x + 0.25, y + 0.22, w - 0.5, 0.55, title, 18, SJTU_RED, True)
    add_textbox(
        slide, x + 0.25, y + 0.98, w - 0.5, h - 1.12, body, 16, DARK,
        line_spacing=1.24
    )


def add_circle_label(slide, x, y, d, text, fill=SJTU_RED, color=WHITE, size=16, bold=True):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(x), Cm(y), Cm(d), Cm(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    add_textbox(
        slide, x, y + d * 0.26, d, d * 0.42, text, size, color, bold,
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.0
    )
    return shape


def add_rule(slide, x, y, w, h=0.035, color=LINE):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_visual_node(slide, x, y, number, title, body, w=6.4):
    add_circle_label(slide, x, y + 0.06, 0.82, number, fill=SJTU_RED, size=14)
    add_textbox(slide, x + 1.05, y, w - 1.05, 0.42, title, 16, SJTU_RED, True)
    add_textbox(slide, x + 1.05, y + 0.52, w - 1.05, 0.52, body, 14, DARK, False, line_spacing=1.2)


def add_card_lines(slide, x, y, w, h, title, lines, accent=False, body_size=16):
    add_card(slide, x, y, w, h, title, "\n".join(lines), accent=accent)


def add_body_slide(prs, title, page_no, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title_bar(slide, title, page_no)
    if subtitle:
        add_textbox(slide, CONTENT_X, 2.65, CONTENT_W, 0.75, subtitle, 20, DARK, True)
    return slide


def add_bullet_list(slide, x, y, w, h, bullets, size=18):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Cm(0.12)
    tf.margin_right = Cm(0.08)
    tf.margin_top = Cm(0.08)
    tf.margin_bottom = Cm(0.08)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        p.level = 0
        p.line_spacing = 1.28
        run = p.add_run()
        run.text = f"• {bullet}"
        set_font(run, size, DARK)


def add_table_cell(slide, x, y, w, h, text, size=12.5, color=DARK, bold=False,
                   fill=None, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill if fill is not None else WHITE
    shape.line.color.rgb = LINE
    add_textbox(
        slide, x + 0.14, y + 0.08, w - 0.28, h - 0.16, text, size, color, bold,
        align=align, valign=MSO_ANCHOR.MIDDLE, margin=0.02, line_spacing=1.14
    )
    return shape


def notes(slide, text):
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = text


def build():
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)

    # Slide 1: cover
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(4.4), Cm(5.05), Cm(25.1), Cm(5.15)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = SJTU_RED
    panel.line.fill.background()
    add_textbox(slide, 5.15, 5.55, 23.6, 1.15, "智能科研方法", 36, WHITE, True,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, 5.15, 6.9, 23.6, 0.72, "AI for Research: Methods and Practice", 20,
                WHITE, False, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, 5.15, 7.85, 23.6, 0.72,
                "主讲人：待补充 · 上海交通大学人工智能研究院", 18,
                WHITE, False, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, 5.15, 8.78, 23.6, 0.58, "开课评审试讲 · 20 分钟", 16,
                WHITE, False, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide, 1)
    notes(slide, "开场先给出课程名称、面向对象和试讲场景。这一页只负责建立课程身份，具体理念放到下一页展开。")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title_bar(slide, "课程简介与设计理念", 2)
    add_textbox(slide, CONTENT_X, 2.65, CONTENT_W, 0.75,
                "用 AI 增强科研全过程，而不是用 AI 替代科研判断。", 20, DARK, True)
    overview = [
        ("课程定位", "面向 CS/AI 研究生\n2 学分 · 32 学时\n专业选修课"),
        ("核心理念", "验证 AI 输出\n记录 AI 参与\n控制科研边界"),
        ("培养目标", "形成可验证、可复现\n可追溯的科研流程\n更可靠地做研究"),
    ]
    for i, (title, body) in enumerate(overview):
        x = COL_XS[i]
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(4.25), Cm(COL_W), Cm(5.85)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else RGBColor(255, 255, 255)
        card.line.color.rgb = LINE
        card.line.width = Pt(1)
        add_circle_label(slide, x + 0.45, 4.72, 1.12, f"{i + 1}", fill=SJTU_RED, size=16)
        add_textbox(slide, x + 1.82, 4.82, COL_W - 2.25, 0.48,
                    title, 18, SJTU_RED, True, valign=MSO_ANCHOR.MIDDLE)
        add_rule(slide, x + 0.65, 6.0, COL_W - 1.3, 0.035, RGBColor(228, 166, 174))
        add_textbox(slide, x + 0.65, 6.75, COL_W - 1.3, 1.65,
                    body, 16, DARK, False, align=PP_ALIGN.CENTER,
                    valign=MSO_ANCHOR.MIDDLE, line_spacing=1.28)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(CONTENT_X), Cm(11.35), Cm(CONTENT_W), Cm(1.3))
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    band.line.color.rgb = LINE
    add_textbox(slide, CONTENT_X + 0.45, 11.62, CONTENT_W - 0.9, 0.46,
                "课程训练学生把 AI 纳入科研流程，而不是把科研判断外包给 AI",
                18, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)
    notes(slide, "这门课的核心定位是科研方法课。AI 在课程中不是目的，而是被纳入可验证、可复现、可追溯科研流程的手段。")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title_bar(slide, "AI 正在改变研究生的科研工作方式", 3)
    add_textbox(slide, CONTENT_X, 2.65, CONTENT_W, 0.75,
                "学生已经自然使用 AI，但未必知道如何验证、记录和控制 AI。", 20, DARK, True)
    items = [
        ("文献检索", "从关键词扩展到证据筛选"),
        ("论文阅读", "贡献、证据、局限结构化"),
        ("代码生成", "辅助实现，但不能污染实验"),
        ("实验迭代", "队列、参数、失败归因"),
        ("结果分析", "从结果到研究判断"),
        ("论文写作", "结论必须回到证据链"),
    ]
    ys = [4.05, 7.15]
    for idx, (title, body) in enumerate(items):
        add_card(slide, COL_XS[idx % 3], ys[idx // 3], COL_W, 2.3, title, body, accent=idx % 2 == 1)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(CONTENT_X), Cm(11.2), Cm(CONTENT_W), Cm(2.15))
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    band.line.color.rgb = LINE
    add_textbox(slide, 2.0, 11.72, 29.2, 0.85,
                "会用 AI 不等于会做研究，课程要补的是验证、记录和控制能力。", 22,
                SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)
    notes(slide, "为什么现在需要这门课，是因为学生已经在真实科研流程里使用 AI。问题不在于学生会不会用 AI，而在于他们是否知道如何验证 AI 输出、记录 AI 参与过程，并控制它不能越过科研判断边界。")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    add_title_bar(slide, "国际课程对标：本课补科研链路", 4)
    add_textbox(slide, CONTENT_X, 2.65, CONTENT_W, 0.75,
                "已有课程证明 Agentic AI 已成体系；本课程借鉴其设计，但落点放在研究生科研方法训练。", 20, DARK, True)
    headers = ["课程", "机构", "特色", "与本课程的关系"]
    col_ws = [7.2, 4.0, 12.2, 7.2]
    table_x = CONTENT_X
    table_y = 3.85
    header_h = 0.72
    row_h = 1.32
    rows = [
        ("CS294/194-196\nLLM Agents", "UC Berkeley",
         "MOOC；25,000+ 学习者；12 讲覆盖推理、规划、代码、多代理、科学发现",
         "理论深度参考；\n进阶路径设计"),
        ("Agentic AI\nCertificate", "Johns Hopkins",
         "16 周在线证书；Python→LLM→Agent→项目；感知-规划-学习-行动框架",
         "渐进式设计参考；\n能力框架互补"),
        ("Agentic AI\nArchitecture", "eCornell",
         "模块化证书；非开发者友好；深度使用 Claude Code 和 LangGraph",
         "模块化设计参考；\n降低门槛策略"),
        ("MAS.664 AI Agents\nand Agentic Web", "MIT",
         "项目驱动；身份、信任、声誉、支付、评估、协调；Internet of Agents",
         "前瞻性主题参考；\n项目制模式"),
    ]
    x = table_x
    for header, w in zip(headers, col_ws):
        add_table_cell(slide, x, table_y, w, header_h, header, 13.5, WHITE, True, SJTU_RED,
                       align=PP_ALIGN.CENTER)
        x += w
    for r, row in enumerate(rows):
        x = table_x
        fill = LIGHT_BG if r % 2 == 0 else WHITE
        for c, (text, w) in enumerate(zip(row, col_ws)):
            add_table_cell(slide, x, table_y + header_h + r * row_h, w, row_h,
                           text, 12.2 if c == 2 else 12.6,
                           SJTU_RED if c == 0 else DARK, c == 0, fill,
                           align=PP_ALIGN.CENTER if c in (0, 1) else PP_ALIGN.LEFT)
            x += w

    add_textbox(slide, CONTENT_X, 10.35, CONTENT_W, 0.52, "对标后的课程定位", 18, SJTU_RED, True)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(CONTENT_X), Cm(11.05), Cm(CONTENT_W), Cm(1.65))
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    band.line.color.rgb = LINE
    add_textbox(slide, CONTENT_X + 0.75, 11.4, CONTENT_W - 1.5, 0.54,
                "不是重复 Agent 工具训练，而是把 AI 纳入可审查、可追溯、负责任的科研链路。",
                20, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)
    notes(slide, "这一页用四门国际课程做对标：Berkeley 提供理论深度和进阶路径参考，Johns Hopkins 提供渐进式能力培养参考，eCornell 提供模块化和降低门槛策略参考，MIT MAS.664 提供前瞻主题和项目制参考。本课程不是简单重复 Agent 工具训练，而是把这些能力纳入研究生科研链路，强调证据链、判断链、留痕链和伦理披露。")

    # Slide 5
    slide = add_body_slide(
        prs, "AI 增强的科研闭环训练", 4,
        "课程按科研动作组织，而不是按工具清单组织。"
    )
    closed_loop_cards = [
        ("课程内容比例", "通用科研工作流  70%-80%\n自动化 / Agent  20%-30%"),
        ("核心原则", "工具服务科研动作\n方法论主线优先\n自动化作为强化模块"),
        ("课程判断", "不是“Agent 自动做任务”\n而是“AI 支撑科研链路”"),
    ]
    for i, (title, body) in enumerate(closed_loop_cards):
        x = CONTENT_X + i * 10.25
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(4.05), Cm(9.45), Cm(5.6)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else RGBColor(255, 255, 255)
        card.line.color.rgb = LINE
        card.line.width = Pt(1)
        add_textbox(slide, x + 0.45, 4.55, 8.55, 0.52,
                    title, 18, SJTU_RED, True)
        add_rule(slide, x + 0.45, 5.45, 8.55, 0.035, RGBColor(228, 166, 174))
        add_textbox(slide, x + 0.55, 6.05, 8.35, 2.85,
                    body, 16, DARK, False, line_spacing=1.18)
    actions = ["找证据", "定问题", "设假设", "做实验", "写论证"]
    for i, action in enumerate(actions):
        x = CONTENT_X + 0.65 + i * 5.65
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(10.9), Cm(4.85), Cm(1.45))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = LINE
        add_textbox(slide, x, 11.18, 4.85, 0.55, action, 18, SJTU_RED, True,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)
    notes(slide, "这门课的组织原则不是工具清单，而是科研动作。大约七到八成内容是通用科研工作流，二到三成是机器学习实验自动化和 Research Agent 实践。每一个工具都必须回答它服务哪个科研动作。")

    # Slide 5
    slide = add_body_slide(
        prs, "八阶段研究链路", 5,
        "AI 可以参与每个阶段，但不能替代关键判断。"
    )
    groups = [
        ("定题", "01 问题定义", "02 第一性原理", "看清问题边界"),
        ("建模", "03 机制假设", "04 外部输入", "形成可检验假设"),
        ("取证", "05 证据整理", "06 研究判断", "把证据转成判断"),
        ("交付", "07 原型验证", "08 回写表达", "沉淀可追溯成果"),
    ]
    block_w = 7.0
    gap = 0.72
    block_y = 4.2
    block_h = 5.9
    for i, (phase, first, second, caption) in enumerate(groups):
        x = CONTENT_X + 0.4 + i * (block_w + gap)
        block = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(block_y), Cm(block_w), Cm(block_h)
        )
        block.fill.solid()
        block.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else RGBColor(255, 255, 255)
        block.line.color.rgb = LINE
        block.line.width = Pt(1)

        add_circle_label(slide, x + 0.35, block_y + 0.34, 1.05, f"{i + 1}", fill=SJTU_RED, size=16)
        add_textbox(slide, x + 1.65, block_y + 0.38, block_w - 2.0, 0.5,
                    phase, 20, SJTU_RED, True, valign=MSO_ANCHOR.MIDDLE)
        add_rule(slide, x + 0.55, block_y + 1.55, block_w - 1.1, 0.035, RGBColor(228, 166, 174))

        for j, stage in enumerate((first, second)):
            pill_y = block_y + 2.05 + j * 1.38
            pill = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x + 0.55), Cm(pill_y), Cm(block_w - 1.1), Cm(0.92)
            )
            pill.fill.solid()
            pill.fill.fore_color.rgb = RGBColor(255, 255, 255)
            pill.line.color.rgb = RGBColor(228, 166, 174)
            add_textbox(slide, x + 0.8, pill_y + 0.15, block_w - 1.6, 0.38,
                        stage, 16, DARK, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

        add_textbox(slide, x + 0.55, block_y + 4.95, block_w - 1.1, 0.45,
                    caption, 14, MUTED, False, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(groups) - 1:
            add_textbox(slide, x + block_w + 0.08, block_y + 2.58, 0.58, 0.55,
                        "→", 20, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(CONTENT_X), Cm(11.05), Cm(CONTENT_W), Cm(1.45))
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    band.line.color.rgb = LINE
    add_textbox(slide, CONTENT_X + 0.45, 11.36, CONTENT_W - 0.9, 0.48,
                "四个不可外包的判断：问题重要性 · 假设可检验性 · 实验公平性 · 结论证据性",
                18, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)
    notes(slide, "课程的主线是八阶段研究链路，从问题定义开始，到第一性原理、机制假设、外部输入摄取、证据整理、研究判断、原型验证，最后回写与表达。AI 可以参与每个阶段，但关键研究判断仍然由人负责。")

    # Slide 6
    slide = add_body_slide(
        prs, "Transformer 论文拆科研闭环", 6,
        "3 分钟微型试讲：不讲模型细节，而是拆问题、假设、证据、判断和留痕。"
    )
    chain = [
        ("问题", "必须依赖\n循环 / 卷积？", "问题定义"),
        ("假设", "注意力能建模\n长距离依赖", "机制假设"),
        ("证据", "翻译结果\n训练效率", "证据地图"),
        ("判断", "结构假设被\n证据支撑", "判断记录"),
        ("风险", "解释性\n成本 / 泛化", "缺口清单"),
    ]
    card_w = 5.45
    gap = 0.58
    top_y = 4.15
    for i, (title, body, artifact) in enumerate(chain):
        x = CONTENT_X + 0.35 + i * (card_w + gap)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(top_y), Cm(card_w), Cm(5.95)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else RGBColor(255, 255, 255)
        card.line.color.rgb = LINE
        card.line.width = Pt(1)
        add_circle_label(slide, x + 0.35, top_y + 0.42, 1.1, f"{i + 1}", fill=SJTU_RED, size=16)
        add_textbox(slide, x + 1.7, top_y + 0.45, card_w - 2.05, 0.48,
                    title, 18, SJTU_RED, True, valign=MSO_ANCHOR.MIDDLE)
        add_rule(slide, x + 0.65, top_y + 1.65, card_w - 1.3, 0.035, RGBColor(228, 166, 174))
        add_textbox(slide, x + 0.45, top_y + 2.25, card_w - 0.9, 1.25,
                    body, 16, DARK, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
        add_textbox(slide, x + 0.5, top_y + 4.45, card_w - 1.0, 0.45,
                    artifact, 14, MUTED, False, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(chain) - 1:
            add_textbox(slide, x + card_w + 0.1, top_y + 2.78, 0.42, 0.48,
                        "→", 18, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(CONTENT_X), Cm(11.55), Cm(CONTENT_W), Cm(1.2))
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    band.line.color.rgb = LINE
    add_textbox(slide, CONTENT_X + 0.4, 11.78, CONTENT_W - 0.8, 0.52,
                "学生不是总结论文，而是抽取证据链、判断链和留痕材料",
                18, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)
    notes(slide, "这一页是试讲片段。我会用 Transformer 经典论文做三分钟微型课堂，但不讲模型细节，而是拆科研链路。学生要学习的不是总结论文，而是抽取问题、假设、证据、判断和留痕材料。")

    # Slide 7
    slide = add_body_slide(
        prs, "32 学时如何组织", 7,
        "六个模块共同服务一个贯穿项目，避免每周孤立练习。"
    )
    modules = [
        ("01", "1-2 周", "导论伦理", "研究工件"),
        ("02", "3-5 周", "文献证据", "研究空白"),
        ("03", "6-7 周", "问题假设", "实验设计"),
        ("04", "8-9 周", "数据基线", "可复现实验"),
        ("05", "10-13 周", "原型验证", "实验自动化"),
        ("06", "14-16 周", "结果分析", "论文表达"),
    ]
    card_w = 4.65
    gap = 0.38
    top_y = 4.25
    for i, (num, week, a, b) in enumerate(modules):
        x = CONTENT_X + 0.35 + i * (card_w + gap)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(top_y), Cm(card_w), Cm(5.3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else RGBColor(255, 255, 255)
        card.line.color.rgb = LINE
        card.line.width = Pt(1)
        add_textbox(slide, x + 0.35, top_y + 0.35, 1.15, 0.55,
                    num, 20, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x + 1.55, top_y + 0.43, card_w - 1.9, 0.42,
                    week, 14, MUTED, True, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)
        add_rule(slide, x + 0.55, top_y + 1.35, card_w - 1.1, 0.035, RGBColor(228, 166, 174))
        add_textbox(slide, x + 0.45, top_y + 2.05, card_w - 0.9, 0.55,
                    a, 16, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x + 0.45, top_y + 3.1, card_w - 0.9, 0.55,
                    b, 16, DARK, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(modules) - 1:
            add_textbox(slide, x + card_w + 0.02, top_y + 2.55, 0.34, 0.48,
                        "→", 16, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    gates = ["问题门", "判断门", "验证门", "论证门"]
    for i, gate in enumerate(gates):
        x = CONTENT_X + 5.2 + i * 5.1
        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(10.15), Cm(3.75), Cm(0.82)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = RGBColor(255, 255, 255)
        pill.line.color.rgb = RGBColor(228, 166, 174)
        add_textbox(slide, x + 0.15, 10.32, 3.45, 0.3,
                    gate, 14, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(CONTENT_X), Cm(11.35), Cm(CONTENT_W), Cm(1.3))
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    band.line.color.rgb = LINE
    add_textbox(slide, CONTENT_X + 0.55, 11.62, CONTENT_W - 1.1, 0.46,
                "一条贯穿项目线：从选题方向到问题门、判断门、验证门、论证门逐步推进",
                18, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)
    notes(slide, "三十二学时会组织成六个模块。所有模块都服务同一个贯穿项目，而不是每周做孤立练习。")

    # Slide 8
    slide = add_body_slide(
        prs, "基础要求与资源边界", 8,
        "课程训练的是方法，不把商业工具购买能力作为学习门槛。"
    )
    resource_cards = [
        ("学生需要", "Python / 机器学习基础\nGit 与命令行\n英文论文阅读"),
        ("学生不需要", "本地 GPU\n自费 API token\n绑定商业工具账号"),
        ("资源口径", "学校已部署资源\n开源或国内可用资源\n免费 token 仅作补充"),
    ]
    for i, (title, body) in enumerate(resource_cards):
        x = CONTENT_X + i * 10.25
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(4.05), Cm(9.45), Cm(6.25)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else RGBColor(255, 255, 255)
        card.line.color.rgb = LINE
        card.line.width = Pt(1)
        add_textbox(slide, x + 0.5, 4.55, 8.45, 0.52,
                    title, 18, SJTU_RED, True)
        add_rule(slide, x + 0.5, 5.48, 8.45, 0.035, RGBColor(228, 166, 174))
        add_textbox(slide, x + 0.6, 6.15, 8.25, 2.55,
                    body, 16, DARK, False, line_spacing=1.32)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(CONTENT_X), Cm(11.4), Cm(CONTENT_W), Cm(1.25))
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    band.line.color.rgb = LINE
    add_textbox(slide, CONTENT_X + 0.4, 11.66, CONTENT_W - 0.8, 0.5,
                "课程可持续运行的前提：资源可替代、权限可控、费用不转嫁给学生",
                18, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)
    notes(slide, "这门课的先修要求是合理但不拔高的。课程资源优先使用学校已部署、开源、国内可用或课程统一配置的模型资源，厂商免费 token 或教育计划只作为补充。")

    # Slide 9
    slide = add_body_slide(
        prs, "课程如何评价学生", 9,
        "评分不是“用了多少 AI 工具”，而是是否形成可验证、可复现、可追溯的科研闭环。"
    )
    score_items = [
        ("25%", "文献与问题定位"),
        ("30%", "实验设计与可复现性"),
        ("25%", "原型与自动化实践"),
        ("20%", "表达、伦理与复盘"),
    ]
    for i, (pct, label) in enumerate(score_items):
        x = CONTENT_X + (i % 2) * 7.4
        y = 4.25 + (i // 2) * 2.6
        add_card(slide, x, y, 6.55, 2.0, pct, label, accent=(i % 2 == 0))
    deliverables = [
        "问题定义", "研究计划", "证据地图", "实验规格",
        "可运行原型", "评价报告", "AI 使用记录", "伦理合规说明",
    ]
    add_textbox(slide, 16.7, 4.15, 13.6, 0.55, "最终提交物：每项都可检查", 18, SJTU_RED, True)
    for i, item in enumerate(deliverables):
        x = 16.7 + (i % 2) * 6.85
        y = 5.15 + (i // 2) * 1.35
        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(6.25), Cm(0.88)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else RGBColor(255, 255, 255)
        pill.line.color.rgb = LINE
        add_textbox(slide, x + 0.22, y + 0.18, 5.8, 0.32,
                    item, 14, DARK, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(CONTENT_X), Cm(11.45), Cm(CONTENT_W), Cm(1.15))
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    band.line.color.rgb = LINE
    add_textbox(slide, CONTENT_X + 0.4, 11.68, CONTENT_W - 0.8, 0.42,
                "每个评分项都有产物支撑，降低“AI 课不好考核”的担忧",
                18, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)
    notes(slide, "考核也围绕科研闭环，而不是围绕用了多少 AI 工具。每一项都有可检查提交物，包括问题定义、研究计划、证据地图、实验规格、可运行原型、评价报告、AI 使用记录、科研伦理与合规说明和论文式短文。")

    # Slide 10
    slide = add_body_slide(
        prs, "研究门与全过程留痕", 10,
        "研究门不是形式化检查，而是防止最后交一份“看似完整但不可验证”的报告。"
    )
    gates = [
        ("问题门", "问题定义", "非例 / 可证伪命题\n证据地图初版"),
        ("判断门", "研究判断", "baseline / 实验规格\n取舍理由"),
        ("验证门", "原型验证", "实验日志 / 失败分析\n评价报告"),
        ("论证门", "结论表达", "追溯链 / AI 披露\n伦理说明"),
    ]
    card_w = 6.95
    gap = 0.72
    top_y = 4.2
    for i, (gate, title, body) in enumerate(gates):
        x = CONTENT_X + 0.25 + i * (card_w + gap)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(top_y), Cm(card_w), Cm(6.05)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else RGBColor(255, 255, 255)
        card.line.color.rgb = LINE
        card.line.width = Pt(1)
        add_circle_label(slide, x + 0.45, top_y + 0.45, 1.18, f"{i + 1}", fill=SJTU_RED, size=16)
        add_textbox(slide, x + 1.85, top_y + 0.52, card_w - 2.25, 0.46,
                    gate, 18, SJTU_RED, True, valign=MSO_ANCHOR.MIDDLE)
        add_rule(slide, x + 0.65, top_y + 1.75, card_w - 1.3, 0.035, RGBColor(228, 166, 174))
        add_textbox(slide, x + 0.65, top_y + 2.45, card_w - 1.3, 0.52,
                    title, 16, DARK, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x + 0.55, top_y + 3.55, card_w - 1.1, 1.15,
                    body, 14, MUTED, False, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.28)
        if i < len(gates) - 1:
            add_textbox(slide, x + card_w + 0.12, top_y + 2.92, 0.42, 0.48,
                        "→", 18, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(CONTENT_X), Cm(11.45), Cm(CONTENT_W), Cm(1.25))
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    band.line.color.rgb = LINE
    add_textbox(slide, CONTENT_X + 0.4, 11.72, CONTENT_W - 0.8, 0.46,
                "每个门都要求留下可审查材料，让 AI 参与过程可复盘、可追责",
                18, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)
    notes(slide, "为了防止项目最后变成一份看似完整但不可验证的报告，课程设置四个研究门。每个门都要求留下可审查材料，让 AI 参与过程可复盘、可追责。")

    # Slide 11
    slide = add_body_slide(
        prs, "主讲人的适配性", 11,
        "保留真实经历待填槽位；不虚构个人履历，不使用空泛形容词。"
    )
    add_card(slide, COL_XS[0], 4.0, COL_W, 6.9, "科研方法积累",
             "问题定义、文献分析、实验设计、论文表达。\n\n待填真实经历：\n研究方向 / 论文 / 项目 / 指导经历", True)
    add_card(slide, COL_XS[1], 4.0, COL_W, 6.9, "AI 与 Agent 实践",
             "AI 辅助研究、代码生成、实验迭代、Research Agent 工作流。\n\n待填真实经历：\n工具链 / Agent 实践 / 自动化实验案例", False)
    add_card(slide, COL_XS[2], 4.0, COL_W, 6.9, "课程转化能力",
             "把个人经验转化为模板、作业、门条件和可评分产出。\n\n待填真实经历：\n课程 / 讲座 / 训练营 / 学生指导经历", True)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(CONTENT_X), Cm(11.55), Cm(CONTENT_W), Cm(1.0))
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    band.line.color.rgb = LINE
    add_textbox(slide, CONTENT_X + 0.4, 11.72, CONTENT_W - 0.8, 0.42,
                "要证明的不是“我熟悉 AI”，而是能把 AI 实践转化为可学习、可操作、可评价的课程",
                16, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)
    notes(slide, "这一页需要替换为主讲人的真实经历。我建议按三条证据组织：科研方法积累、AI 与 Agent 实践、课程转化能力。这里宁可留待补充，也不虚构履历。")

    # Slide 12
    slide = add_body_slide(
        prs, "课程产出：研究雏形", 12,
        "学生拿走的不是工具清单，而是一套能延续到自己科研工作的材料和习惯。"
    )
    core = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(CONTENT_X + 10.55), Cm(5.2), Cm(9.5), Cm(3.3)
    )
    core.fill.solid()
    core.fill.fore_color.rgb = SJTU_RED
    core.line.color.rgb = SJTU_RED
    add_textbox(slide, CONTENT_X + 11.15, 5.82, 8.3, 0.55,
                "可延续的研究雏形", 20, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, CONTENT_X + 11.25, 6.75, 8.1, 0.6,
                "问题 · 证据 · 实验 · 表达", 16, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    artifacts = [
        (CONTENT_X + 0.4, 4.25, "问题定义", "研究计划"),
        (CONTENT_X + 0.4, 8.2, "证据地图", "文献与空白"),
        (CONTENT_X + 21.2, 4.25, "实验规格", "baseline / 日志"),
        (CONTENT_X + 21.2, 8.2, "评价报告", "论文式短文"),
    ]
    for i, (x, y, title, body) in enumerate(artifacts):
        add_card(slide, x, y, 8.2, 2.75, title, body, accent=(i % 2 == 0))
        if x < CONTENT_X + 10:
            add_textbox(slide, x + 8.3, y + 1.1, 0.75, 0.45, "→", 18, SJTU_RED, True,
                        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        else:
            add_textbox(slide, x - 0.75, y + 1.1, 0.75, 0.45, "←", 18, SJTU_RED, True,
                        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    support = [
        ("AI 使用记录", CONTENT_X + 9.3, 9.25),
        ("伦理合规说明", CONTENT_X + 15.65, 9.25),
    ]
    for text, x, y in support:
        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(5.35), Cm(0.92)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = RGBColor(255, 255, 255)
        pill.line.color.rgb = RGBColor(228, 166, 174)
        add_textbox(slide, x + 0.2, y + 0.18, 4.95, 0.32,
                    text, 14, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(CONTENT_X), Cm(11.3), Cm(CONTENT_W), Cm(1.25))
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    band.line.color.rgb = LINE
    add_textbox(slide, CONTENT_X + 0.4, 11.55, CONTENT_W - 0.8, 0.5,
                "目标不是“泛化提升科研能力”，而是产出可延续的研究材料",
                18, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)
    notes(slide, "学生学完这门课，应该得到一个可继续发展的研究雏形。这些产出可以继续服务开题准备、组会汇报、论文 idea 打磨、baseline 建设和后续科研项目启动。")

    # Slide 13
    slide = add_body_slide(
        prs, "科研过程必须可审查", 13,
        "课程不是禁止学生用 AI，而是要求 AI 参与过程留下证据和边界。"
    )
    columns = [
        ("输入可追溯", "AI 使用记录", "提示词 / 来源 / 工具\n人工核验"),
        ("权限可说明", "工具权限记录", "数据源 / 读写范围\n外部服务"),
        ("实验可复现", "实验留痕", "实验 ID / 配置\n代码 / 结果"),
        ("判断可问责", "研究判断留痕", "证据地图 / 判断记录\n评价报告"),
    ]
    top_y = 4.25
    card_w = 7.05
    gap = 0.55
    for i, (label, title, body) in enumerate(columns):
        x = CONTENT_X + 0.2 + i * (card_w + gap)
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(top_y), Cm(card_w), Cm(5.95)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else RGBColor(255, 255, 255)
        card.line.color.rgb = LINE
        card.line.width = Pt(1)
        add_circle_label(slide, x + 0.45, top_y + 0.45, 1.25, f"{i + 1}", fill=SJTU_RED, size=16)
        add_textbox(slide, x + 2.0, top_y + 0.45, card_w - 2.4, 0.48,
                    label, 16, MUTED, True, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x + 0.55, top_y + 2.05, card_w - 1.1, 0.58,
                    title, 18, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_rule(slide, x + 0.85, top_y + 3.2, card_w - 1.7, 0.035, RGBColor(228, 166, 174))
        add_textbox(slide, x + 0.75, top_y + 3.65, card_w - 1.5, 1.35,
                    body, 14, DARK, False, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=1.28)
        if i < len(columns) - 1:
            add_textbox(slide, x + card_w + 0.08, top_y + 2.62, 0.4, 0.48,
                        "→", 18, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    boundary = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(CONTENT_X + 7.4), Cm(11.0), Cm(16.1), Cm(1.25)
    )
    boundary.fill.solid()
    boundary.fill.fore_color.rgb = RGBColor(255, 255, 255)
    boundary.line.color.rgb = RGBColor(228, 166, 174)
    add_textbox(slide, CONTENT_X + 7.85, 11.27, 15.2, 0.42,
                "伦理边界：数据 · 隐私 · 许可 · 署名 · AI 使用披露",
                16, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(CONTENT_X), Cm(13.0), Cm(CONTENT_W), Cm(1.25))
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_BG
    band.line.color.rgb = LINE
    add_textbox(slide, CONTENT_X + 0.4, 13.32, CONTENT_W - 0.8, 0.46,
                "不绑定单一商业工具 · 不要求学生自费购买 token · 免费 token 或教育计划只作为补充资源",
                16, SJTU_RED, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide)
    notes(slide, "课程不是禁止学生使用 AI，而是要求 AI 参与过程可审查、可复盘、可问责。AI 使用记录、工具权限记录、实验留痕、研究判断留痕和伦理边界都要留下证据。")

    # Slide 14
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    white_bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(33.87), Cm(19.05)
    )
    white_bg.fill.solid()
    white_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    white_bg.line.fill.background()
    add_textbox(slide, 3.0, 3.25, 27.8, 0.78,
                "智能科研方法", 26, SJTU_RED, True,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, 4.0, 4.55, 25.8, 0.6,
                "AI for Research: Methods and Practice", 16, MUTED, False,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, 4.0, 6.5, 25.8, 1.55,
                "培养能在 AI 时代更可靠地做研究的研究生", 22, DARK, True,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, 5.2, 8.6, 23.4, 1.15,
                "把 AI 纳入可验证、可复现、可追溯的科研流程", 18, DARK, False,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    red_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(12.55), Cm(33.87), Cm(4.35)
    )
    red_band.fill.solid()
    red_band.fill.fore_color.rgb = SJTU_RED
    red_band.line.fill.background()
    add_textbox(slide, 4.0, 13.85, 25.8, 0.9,
                "谢谢批评指正", 32, WHITE, True,
                align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    notes(slide, "最后回到课程价值。这门课不是教学生把研究外包给 AI，而是教学生把 AI 纳入可验证、可复现、可追溯的科研流程。最终目标是培养能在 AI 时代更可靠地做研究的研究生。")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
